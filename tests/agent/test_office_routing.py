"""Document-deliverable routing: intent detection + inline-build backstop."""

from __future__ import annotations

import importlib

office_routing = importlib.import_module("agent.office_routing")
detect = office_routing.detect_document_deliverable_intent
is_inline_doc_build = office_routing.is_inline_doc_build


class TestIntentDetection:
    def test_fires_on_real_requests(self):
        for msg in [
            "make me a 6-slide deck on the compression thesis",
            "i want you to made a docx word doc with the original story",
            "create a Word document summarizing this",
            "build a PowerPoint presentation about our results",
            "can you put together a report on the P2 data",
            "generate a pptx from these notes",
            "draft a proposal for the client",
            "turn this into a polished slide deck",
        ]:
            assert detect(msg), f"should detect: {msg!r}"

    def test_ignores_non_deliverables(self):
        for msg in [
            "document this function with a docstring",  # bare 'document', no deliverable noun
            "what does this report say?",               # no creation verb
            "summarize the attached file",
            "read the slides and tell me the gist",     # read, not create
            "fix the bug in build_docx.py",
            "",
            None,
        ]:
            assert not detect(msg), f"should NOT detect: {msg!r}"

    def test_disabled_by_env(self, monkeypatch):
        monkeypatch.setenv("HERMES_DISABLE_OFFICE_ROUTING", "1")
        assert not detect("make me a deck")


class TestInlineBuildBackstop:
    def test_flags_inline_docx_build(self):
        code = (
            "from docx import Document\n"
            "d = Document()\n"
            "d.add_heading('Title')\n"
            "d.add_paragraph('body')\n"
            "d.save('/tmp/out.docx')\n"
        )
        assert is_inline_doc_build(code)

    def test_flags_inline_pptx_build(self):
        code = (
            "import pptx\n"
            "p = pptx.Presentation()\n"
            "p.slides.add_slide(p.slide_layouts[0])\n"
            "p.save('/tmp/out.pptx')\n"
        )
        assert is_inline_doc_build(code)

    def test_allows_matplotlib_illustration(self):
        # Generating an illustration IMAGE is legitimate (embedded via an image
        # block; the render_check gate verifies the final pages). Only inline
        # python-docx/pptx DOCUMENT building is blocked.
        code = "import matplotlib.pyplot as plt\nplt.plot([1,2])\nplt.savefig('/tmp/x.png')\n"
        assert not is_inline_doc_build(code)

    def test_allows_reading_existing_docx(self):
        # Reading (no .save) must not be blocked — non-deliverable doc work.
        code = "from docx import Document\nd = Document('/tmp/in.docx')\nprint(d.paragraphs[0].text)\n"
        assert not is_inline_doc_build(code)

    def test_allows_unrelated_code(self):
        assert not is_inline_doc_build("import os\nprint(os.getcwd())\n")
        assert not is_inline_doc_build("")


def test_contextvar_roundtrip(monkeypatch):
    monkeypatch.delenv("HERMES_DISABLE_OFFICE_ROUTING", raising=False)
    office_routing.set_office_deliverable_turn(True)
    assert office_routing.is_office_deliverable_turn() is True
    office_routing.set_office_deliverable_turn(False)
    assert office_routing.is_office_deliverable_turn() is False
    # Disabled env forces the flag off even when set True.
    monkeypatch.setenv("HERMES_DISABLE_OFFICE_ROUTING", "1")
    office_routing.set_office_deliverable_turn(True)
    assert office_routing.is_office_deliverable_turn() is False
