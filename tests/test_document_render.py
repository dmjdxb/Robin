"""M0 render-layer tests: rasterise documents to per-page images for the vision QA gate.

PDF rendering is hermetic (PyMuPDF generates the fixture and rasterises it, no system
deps). LibreOffice-dependent office rendering is exercised only when soffice is present.
"""
import os

import pytest

pymupdf = pytest.importorskip("pymupdf")

from tools.document_render import (  # noqa: E402
    find_soffice,
    render_to_images,
    soffice_available,
)


def _make_pdf(path: str, pages: int = 3, width: int = 720, height: int = 405) -> None:
    doc = pymupdf.open()
    for n in range(pages):
        page = doc.new_page(width=width, height=height)
        page.insert_text((40, 60), f"Slide {n + 1}", fontsize=24)
    doc.save(path)
    doc.close()


def test_render_pdf_to_page_images(tmp_path):
    pdf = tmp_path / "deck.pdf"
    _make_pdf(str(pdf), pages=3)
    imgs = render_to_images(str(pdf), outdir=str(tmp_path / "imgs"))
    assert len(imgs) == 3
    for p in imgs:
        assert p.endswith(".png")
        assert os.path.getsize(p) > 0


def test_longest_side_capped(tmp_path):
    pdf = tmp_path / "big.pdf"
    doc = pymupdf.open()
    doc.new_page(width=1920, height=1080)  # large page
    doc.save(str(pdf))
    doc.close()
    imgs = render_to_images(str(pdf), outdir=str(tmp_path / "imgs"), max_px=1024)
    pix = pymupdf.Pixmap(imgs[0])
    assert max(pix.width, pix.height) <= 1024  # capped so vision input stays cheap


def test_find_soffice_returns_path_or_none():
    r = find_soffice()
    assert r is None or isinstance(r, str)
    assert isinstance(soffice_available(), bool)


@pytest.mark.skipif(not soffice_available(), reason="LibreOffice not installed")
def test_render_pptx_when_libreoffice_present(tmp_path):
    pptx = pytest.importorskip("pptx")
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.add_textbox(Inches(1), Inches(1), Inches(6), Inches(1))
    out = tmp_path / "x.pptx"
    prs.save(str(out))
    imgs = render_to_images(str(out), outdir=str(tmp_path / "imgs"))
    assert len(imgs) >= 1
