"""M1 builder tests: the constrained PPTX builder fills designed layouts from a deck spec."""
import pytest

pytest.importorskip("pptx")

from tools.office_pptx import THEMES, build_deck  # noqa: E402

_SPEC = {
    "title": "Quarterly Review",
    "theme": "light",
    "slides": [
        {"layout": "title", "title": "Quarterly Review", "subtitle": "Q2 2026"},
        {"layout": "section", "title": "Results"},
        {"layout": "bullets", "title": "Highlights", "bullets": ["Revenue up", "Churn down", "NPS up"]},
        {
            "layout": "two_column",
            "title": "Wins vs Risks",
            "left_title": "Wins",
            "left": ["Launched X", "Signed Y"],
            "right_title": "Risks",
            "right": ["Hiring", "Latency"],
        },
        {"layout": "quote", "quote": "Make it work, then make it good.", "attribution": "Team"},
        {"layout": "closing", "title": "Thank you", "subtitle": "Questions?"},
    ],
}


def test_build_all_layouts(tmp_path):
    out = tmp_path / "deck.pptx"
    res = build_deck(_SPEC, str(out))
    assert out.exists() and out.stat().st_size > 0
    assert res["slides"] == 6
    assert res["warnings"] == []
    # re-open to confirm it's a valid package with the right slide count
    from pptx import Presentation

    prs = Presentation(str(out))
    assert len(prs.slides) == 6
    # 16:9 canvas
    assert round(prs.slide_width / 914400, 2) == 13.33  # EMUs per inch


def test_unknown_layout_falls_back_with_warning(tmp_path):
    spec = {"slides": [{"layout": "nope", "bullets": ["a"]}]}
    res = build_deck(spec, str(tmp_path / "x.pptx"))
    assert res["slides"] == 1
    assert any("unknown layout" in w for w in res["warnings"])


def test_themes_exist():
    assert "light" in THEMES and "dark" in THEMES


def test_bad_spec_raises(tmp_path):
    with pytest.raises(ValueError):
        build_deck({"no_slides": True}, str(tmp_path / "x.pptx"))


def test_table_chart_image_layouts(tmp_path):
    pymupdf = pytest.importorskip("pymupdf")
    d = pymupdf.open()
    pg = d.new_page(width=200, height=120)
    pg.insert_text((10, 40), "img", fontsize=20)
    img = tmp_path / "pic.png"
    pg.get_pixmap().save(str(img))
    d.close()

    spec = {
        "theme": "light",
        "slides": [
            {"layout": "table", "title": "T", "headers": ["A", "B"], "rows": [["1", "2"], ["3", "4"]]},
            {"layout": "chart", "title": "C", "chart_type": "column",
             "categories": ["Q1", "Q2"], "series": {"Rev": [10, 20]}},
            {"layout": "image", "title": "I", "image_path": str(img), "caption": "cap"},
        ],
    }
    out = tmp_path / "d.pptx"
    res = build_deck(spec, str(out))
    assert res["slides"] == 3 and res["warnings"] == []

    from pptx import Presentation

    prs = Presentation(str(out))
    assert any(getattr(sh, "has_table", False) for sh in prs.slides[0].shapes)
    assert any(getattr(sh, "has_chart", False) for sh in prs.slides[1].shapes)
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    assert any(sh.shape_type == MSO_SHAPE_TYPE.PICTURE for sh in prs.slides[2].shapes)
