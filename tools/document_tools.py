#!/usr/bin/env python3
"""Ask your document — grounded Q&A over office documents (PDF/DOCX/PPTX/XLSX).

This is Robin's flagship "talk to your document" tool. It does NOT answer the
question itself: it extracts the document into clean, citation-anchored text and
returns that text to Robin's main agent, which already holds the conversation and
a long context window. The agent then answers grounded in the returned text and
cites the bracketed anchors (e.g. ``[p.3]``, ``[slide 4]``).

Design (Tier 1 — whole-document, no vector DB):
  * Documents that fit a character budget are returned whole — best coherence.
  * Larger documents are returned as an OUTLINE of stable sections; the agent
    re-calls with ``section="<id>"`` to pull just the relevant part. This keeps
    the document in the normal tool-result lifecycle (so the trajectory
    compressor never silently summarises it away) and never blows the result cap.

Citation anchors per format (human-meaningful — they match what a user sees):
  PDF   -> [p.N]                     DOCX  -> [§ Heading]  (fallback [¶ N])
  PPTX  -> [slide N]                 XLSX  -> [Sheet:Name] + cell-precise [Sheet!C12]

Heavy parsing libraries (pypdf/pdfplumber/python-docx/python-pptx/openpyxl) are
LAZY-installed at first use via tools.lazy_deps — this module imports cleanly
without them present, so a user who never opens a document never pays for them.
"""

from __future__ import annotations

import logging
from pathlib import Path
from typing import Optional

from tools.registry import registry
from tools.file_tools import _resolve_path_for_task

logger = logging.getLogger(__name__)

# ── lazy-bound library handles (None until the format is first used) ──────────
pypdf = None
pdfplumber = None
docx = None        # python-docx imports as `docx`
pptx = None        # python-pptx imports as `pptx`
openpyxl = None

# Return-size budget. The registry cap is 100K chars; stay under it with headroom
# for the framing header/footer. A document whose extracted text exceeds this is
# returned as an OUTLINE instead of whole.
_BUDGET_CHARS = 80_000
# Spreadsheet rendering caps so a giant sheet can't dominate the result.
_XLSX_MAX_ROWS = 200
_XLSX_MAX_COLS = 40

_SUPPORTED = {".pdf", ".docx", ".pptx", ".xlsx"}
# Legacy binary formats we can't read directly — guide the user to re-save.
_LEGACY = {".doc": ".docx", ".ppt": ".pptx", ".xls": ".xlsx"}


# ─────────────────────────────────────────────────────────────────────────────
# Lazy dependency binding (one feature key per format — a DOCX user never pulls
# the PDF libraries).
# ─────────────────────────────────────────────────────────────────────────────
def _ensure_pdf() -> bool:
    global pypdf, pdfplumber
    if pypdf is not None:
        return True

    def _imp() -> dict:
        import pypdf as _pypdf
        try:
            import pdfplumber as _pdfplumber
        except Exception:  # pdfplumber is best-effort table support
            _pdfplumber = None
        return {"pypdf": _pypdf, "pdfplumber": _pdfplumber}

    from tools.lazy_deps import ensure_and_bind
    return ensure_and_bind("document.pdf", _imp, globals(), prompt=False)


def _ensure_docx() -> bool:
    global docx
    if docx is not None:
        return True

    def _imp() -> dict:
        import docx as _docx
        return {"docx": _docx}

    from tools.lazy_deps import ensure_and_bind
    return ensure_and_bind("document.docx", _imp, globals(), prompt=False)


def _ensure_pptx() -> bool:
    global pptx
    if pptx is not None:
        return True

    def _imp() -> dict:
        import pptx as _pptx
        return {"pptx": _pptx}

    from tools.lazy_deps import ensure_and_bind
    return ensure_and_bind("document.pptx", _imp, globals(), prompt=False)


def _ensure_xlsx() -> bool:
    global openpyxl
    if openpyxl is not None:
        return True

    def _imp() -> dict:
        import openpyxl as _openpyxl
        return {"openpyxl": _openpyxl}

    from tools.lazy_deps import ensure_and_bind
    return ensure_and_bind("document.xlsx", _imp, globals(), prompt=False)


# ─────────────────────────────────────────────────────────────────────────────
# Per-format extraction. Each returns (chunks, meta) where chunks is an ordered
# list of {"anchor": str, "text": str} and meta describes the document.
# ─────────────────────────────────────────────────────────────────────────────
def _extract_pdf(path: Path) -> tuple[list[dict], dict]:
    reader = pypdf.PdfReader(str(path))
    pages = reader.pages
    chunks: list[dict] = []
    for i, page in enumerate(pages, start=1):
        try:
            text = (page.extract_text() or "").strip()
        except Exception:
            text = ""
        chunks.append({"anchor": f"[p.{i}]", "text": text})

    # Best-effort table augmentation: legal/marketing PDFs lean on tables, which
    # pypdf flattens. Append markdown tables (if any) under each page anchor.
    if pdfplumber is not None:
        try:
            with pdfplumber.open(str(path)) as pdf:
                for i, page in enumerate(pdf.pages, start=1):
                    if i > len(chunks):
                        break
                    tables = page.extract_tables() or []
                    md = "\n\n".join(_rows_to_md(t) for t in tables if t)
                    if md.strip():
                        chunks[i - 1]["text"] = (chunks[i - 1]["text"] + "\n\n" + md).strip()
        except Exception:
            logger.debug("pdfplumber table pass failed; using pypdf text only", exc_info=True)

    text_chars = sum(len(c["text"]) for c in chunks)
    meta = {"type": "PDF", "unit": "pages", "count": len(pages),
            "scanned": (len(pages) > 0 and text_chars < 40 * len(pages))}
    return chunks, meta


def _iter_docx_blocks(document):
    """Yield paragraphs and tables in document order (python-docx splits them).

    Detect element type by XML tag (``…}p`` / ``…}tbl``) rather than importing
    version-specific internal classes — keeps this robust across python-docx
    releases.
    """
    from docx.text.paragraph import Paragraph
    from docx.table import Table

    for child in document.element.body.iterchildren():
        tag = getattr(child, "tag", "")
        if tag.endswith("}p"):
            yield Paragraph(child, document)
        elif tag.endswith("}tbl"):
            yield Table(child, document)


def _extract_docx(path: Path) -> tuple[list[dict], dict]:
    from docx.text.paragraph import Paragraph
    from docx.table import Table

    document = docx.Document(str(path))
    chunks: list[dict] = []
    cur_anchor = "[¶ 1]"
    cur_buf: list[str] = []
    para_idx = 0
    headings = 0

    def _flush():
        if cur_buf:
            chunks.append({"anchor": cur_anchor, "text": "\n".join(cur_buf).strip()})

    for block in _iter_docx_blocks(document):
        if isinstance(block, Paragraph):
            style = (block.style.name if block.style else "") or ""
            txt = block.text.strip()
            if style.startswith("Heading") and txt:
                _flush()
                cur_buf = []
                headings += 1
                cur_anchor = f"[§ {txt}]"
            else:
                para_idx += 1
                if txt:
                    cur_buf.append(txt)
        elif isinstance(block, Table):
            rows = [[cell.text for cell in row.cells] for row in block.rows]
            md = _rows_to_md(rows)
            if md.strip():
                cur_buf.append(md)
    _flush()

    if not chunks:
        chunks = [{"anchor": "[¶ 1]", "text": ""}]
    meta = {"type": "Word document", "unit": "sections", "count": len(chunks)}
    return chunks, meta


def _extract_pptx(path: Path) -> tuple[list[dict], dict]:
    prs = pptx.Presentation(str(path))
    chunks: list[dict] = []
    slides = list(prs.slides)
    for i, slide in enumerate(slides, start=1):
        title = ""
        try:
            if slide.shapes.title and slide.shapes.title.text:
                title = slide.shapes.title.text.strip()
        except Exception:
            title = ""
        body: list[str] = []
        for shape in slide.shapes:
            try:
                if shape.has_text_frame:
                    t = shape.text_frame.text.strip()
                    if t and t != title:
                        body.append(t)
                if shape.has_table:
                    rows = [[c.text for c in r.cells] for r in shape.table.rows]
                    md = _rows_to_md(rows)
                    if md.strip():
                        body.append(md)
            except Exception:
                continue
        header = f'[slide {i}] "{title}"' if title else f"[slide {i}]"
        chunks.append({"anchor": header, "text": "\n".join(body).strip()})

        # speaker notes as a separate citable anchor
        try:
            if slide.has_notes_slide:
                notes = (slide.notes_slide.notes_text_frame.text or "").strip()
                if notes:
                    chunks.append({"anchor": f"[slide {i} · notes]", "text": notes})
        except Exception:
            pass

    meta = {"type": "PowerPoint", "unit": "slides", "count": len(slides)}
    return chunks, meta


def _extract_xlsx(path: Path) -> tuple[list[dict], dict]:
    wb = openpyxl.load_workbook(str(path), read_only=True, data_only=True)
    chunks: list[dict] = []
    for ws in wb.worksheets:
        rows = list(ws.iter_rows(values_only=True))
        truncated = False
        if len(rows) > _XLSX_MAX_ROWS:
            rows = rows[:_XLSX_MAX_ROWS]
            truncated = True
        # Render with column letters + row numbers so cells stay citable as
        # [Sheet!C12] (coordinates preserved in the markdown header/index).
        n_cols = min(max((len(r) for r in rows), default=0), _XLSX_MAX_COLS)
        if n_cols == 0:
            chunks.append({"anchor": f"[Sheet:{ws.title}]", "text": "(empty sheet)"})
            continue
        col_letters = [_col_letter(c + 1) for c in range(n_cols)]
        lines = ["| | " + " | ".join(col_letters) + " |",
                 "|---|" + "---|" * n_cols]
        for ridx, row in enumerate(rows, start=1):
            cells = [("" if v is None else str(v)) for v in row[:n_cols]]
            cells += [""] * (n_cols - len(cells))
            lines.append(f"| {ridx} | " + " | ".join(cells) + " |")
        body = "\n".join(lines)
        if truncated:
            body += f"\n(showing first {_XLSX_MAX_ROWS} rows)"
        chunks.append({"anchor": f"[Sheet:{ws.title}]", "text": body})
    wb.close()
    if not chunks:
        chunks = [{"anchor": "[Sheet:1]", "text": "(empty workbook)"}]
    meta = {"type": "Excel spreadsheet", "unit": "sheets", "count": len(chunks)}
    return chunks, meta


# ─────────────────────────────────────────────────────────────────────────────
# Small formatting helpers
# ─────────────────────────────────────────────────────────────────────────────
def _col_letter(n: int) -> str:
    """1 -> A, 27 -> AA (spreadsheet column letters)."""
    s = ""
    while n > 0:
        n, r = divmod(n - 1, 26)
        s = chr(65 + r) + s
    return s


def _rows_to_md(rows: list[list]) -> str:
    """Render a list-of-rows as a GitHub-flavoured markdown table."""
    rows = [["" if c is None else str(c).replace("\n", " ").strip() for c in row] for row in rows if row]
    if not rows:
        return ""
    width = max(len(r) for r in rows)
    rows = [r + [""] * (width - len(r)) for r in rows]
    head, *body = rows
    out = ["| " + " | ".join(head) + " |", "|" + "---|" * width]
    out += ["| " + " | ".join(r) + " |" for r in body]
    return "\n".join(out)


def _render_full(chunks: list[dict], meta: dict, filename: str) -> str:
    head = (
        f"=== DOCUMENT: {filename} ===\n"
        f"Type: {meta['type']} · {meta['count']} {meta['unit']} · native text\n"
        "Cite the bracketed anchor before each block, e.g. [p.3]. "
        "Answer ONLY from the text below; if the answer is not here, say so.\n"
    )
    parts = [head]
    if meta.get("scanned"):
        parts.append(
            "\nNOTE: this document appears to be SCANNED (little or no selectable "
            "text). Reading scanned documents is coming soon; for now, answer only "
            "from whatever text was recovered below.\n"
        )
    for c in chunks:
        body = c["text"] or "(no text on this " + meta["unit"][:-1] + ")"
        parts.append(f"\n{c['anchor']}\n{body}\n")
    parts.append(f"\n=== END DOCUMENT ({filename}) ===")
    return "".join(parts)


def _render_outline(chunks: list[dict], meta: dict, filename: str) -> str:
    """Group oversized documents into budget-sized sections with stable ids."""
    sections: list[tuple[int, int, str, str]] = []  # (start, end, first_anchor, last_anchor)
    start = 0
    size = 0
    for i, c in enumerate(chunks):
        size += len(c["text"]) + len(c["anchor"]) + 2
        is_last = i == len(chunks) - 1
        if size >= _BUDGET_CHARS or is_last:
            sections.append((start, i, chunks[start]["anchor"], chunks[i]["anchor"]))
            start = i + 1
            size = 0

    lines = [
        f"=== DOCUMENT OUTLINE: {filename} ({meta['count']} {meta['unit']}, too large to load whole) ===",
        "This document exceeds the single-load budget, so it is split into sections.",
        'To read a section, call ask_document again with section="<id>".',
        "Sections:",
    ]
    for (s, e, a0, a1) in sections:
        sid = f"{s + 1}-{e + 1}"  # 1-indexed, inclusive
        rng = a0 if a0 == a1 else f"{a0}–{a1}"
        lines.append(f"  - {rng}   (id: {sid})")
    lines.append("Tip: if the question spans the whole document, fetch sections in order.")
    lines.append("=== END OUTLINE ===")
    return "\n".join(lines)


def _slice_section(chunks: list[dict], section: str) -> Optional[list[dict]]:
    """Resolve a section id like '1-20' (1-indexed inclusive) to its chunks."""
    try:
        a, b = section.strip().split("-", 1)
        lo, hi = int(a), int(b)
    except Exception:
        return None
    lo = max(1, lo)
    hi = min(len(chunks), hi)
    if lo > hi:
        return None
    return chunks[lo - 1:hi]


# ─────────────────────────────────────────────────────────────────────────────
# Core entrypoint
# ─────────────────────────────────────────────────────────────────────────────
def ask_document_tool(path: str, section: Optional[str] = None, task_id: str = "default") -> str:
    if not path or not str(path).strip():
        return "ERROR: no document path provided. Pass the path to a PDF, Word, PowerPoint, or Excel file."

    resolved = _resolve_path_for_task(path, task_id)
    if not resolved.exists():
        return f"ERROR: document not found: {path}"
    if resolved.is_dir():
        return f"ERROR: {path} is a folder, not a document."

    suffix = resolved.suffix.lower()
    if suffix in _LEGACY:
        return (
            f"ERROR: '{resolved.name}' is an older binary {suffix} file, which can't be read "
            f"directly. Please re-save it as {_LEGACY[suffix]} (File → Save As) and try again."
        )
    if suffix not in _SUPPORTED:
        return (
            f"ERROR: unsupported document type '{suffix}'. Supported: PDF, Word (.docx), "
            f"PowerPoint (.pptx), Excel (.xlsx). For plain-text files use read_file instead."
        )

    try:
        if suffix == ".pdf":
            if not _ensure_pdf():
                return "ERROR: could not install the PDF reader. Check the network connection and try again."
            chunks, meta = _extract_pdf(resolved)
        elif suffix == ".docx":
            if not _ensure_docx():
                return "ERROR: could not install the Word reader. Check the network connection and try again."
            chunks, meta = _extract_docx(resolved)
        elif suffix == ".pptx":
            if not _ensure_pptx():
                return "ERROR: could not install the PowerPoint reader. Check the network connection and try again."
            chunks, meta = _extract_pptx(resolved)
        else:  # .xlsx
            if not _ensure_xlsx():
                return "ERROR: could not install the Excel reader. Check the network connection and try again."
            chunks, meta = _extract_xlsx(resolved)
    except Exception as e:  # noqa: BLE001 — surface a clean message, never raise to the loop
        logger.exception("ask_document extraction failed for %s", resolved)
        return f"ERROR: could not read '{resolved.name}': {e}"

    # Section fetch (follow-up to an OUTLINE) — return just that slice, whole.
    if section:
        sliced = _slice_section(chunks, section)
        if sliced is None:
            return f"ERROR: unknown section id '{section}'. Call ask_document without a section to see the outline."
        return _render_full(sliced, meta, resolved.name)

    total = sum(len(c["text"]) for c in chunks)
    if total <= _BUDGET_CHARS:
        return _render_full(chunks, meta, resolved.name)
    return _render_outline(chunks, meta, resolved.name)


# ─────────────────────────────────────────────────────────────────────────────
# Tool registration
# ─────────────────────────────────────────────────────────────────────────────
ASK_DOCUMENT_SCHEMA = {
    "name": "ask_document",
    "description": (
        "Read a document the user is asking about (PDF, Word .docx, PowerPoint .pptx, or "
        "Excel .xlsx) and return its text with citation anchors. Use this whenever the user "
        "attaches a document and asks about it. "
        "RULES FOR YOUR ANSWER: (1) Answer ONLY from the returned text — do not add outside "
        "knowledge. (2) Cite the bracketed anchor for every claim, e.g. \"Net-30 applies [p.3]\" — "
        "use the exact anchor shown ([p.N] pages, [slide N] slides, [§ Heading] Word, [Sheet!Cell] "
        "Excel). (3) If the answer is not in the document, say \"I couldn't find that in this "
        "document\" — do not guess. (4) If an OUTLINE is returned, call this tool again with the "
        "relevant section id before answering. (5) The document text stays in the conversation — "
        "do not re-call for the same content you already have."
    ),
    "parameters": {
        "type": "object",
        "properties": {
            "path": {
                "type": "string",
                "description": "Path to the document (PDF/DOCX/PPTX/XLSX); absolute, relative, or ~/path.",
            },
            "section": {
                "type": "string",
                "description": (
                    "Optional. Only when a previous call returned an OUTLINE: pass a section id "
                    "(e.g. '1-20') to fetch that section's full text. Omit on the first call."
                ),
            },
        },
        "required": ["path"],
    },
}


def _check_ask_document_reqs() -> bool:
    """Always offer the tool; parser libraries lazy-install on first use."""
    return True


def _handle_ask_document(args, **kw):
    tid = kw.get("task_id") or "default"
    return ask_document_tool(
        path=args.get("path", ""),
        section=args.get("section") or None,
        task_id=tid,
    )


registry.register(
    name="ask_document",
    toolset="file",
    schema=ASK_DOCUMENT_SCHEMA,
    handler=_handle_ask_document,
    check_fn=_check_ask_document_reqs,
    emoji="📄",
    max_result_size_chars=100_000,
)
