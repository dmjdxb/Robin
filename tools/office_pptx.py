"""Constrained PPTX builder — the layout fix.

Instead of the model guessing slide coordinates (the old blind PptxGenJS path), it
fills CONTENT into a small library of designed, tested layouts defined here. Geometry,
type scale, palette and spacing are fixed once per theme; the model only supplies text
and data via a deck spec. Text frames shrink-to-fit, and the vision QA gate (M2) catches
any residual overflow. This is the "templates so layout is never guessed" pillar.

Deck spec (JSON):
  {
    "title": "...", "theme": "light",
    "slides": [
      {"layout": "title",      "title": "...", "subtitle": "..."},
      {"layout": "section",    "title": "..."},
      {"layout": "bullets",    "title": "...", "bullets": ["...", "..."]},
      {"layout": "two_column", "title": "...",
       "left_title": "...", "left": ["..."], "right_title": "...", "right": ["..."]},
      {"layout": "quote",      "quote": "...", "attribution": "..."},
      {"layout": "closing",    "title": "...", "subtitle": "..."}
    ]
  }
"""
from __future__ import annotations

import json
from dataclasses import dataclass
from typing import Any, Optional

# Lazy-bound python-pptx symbols (installed on first office use via document.pptx).
Presentation = Inches = Pt = RGBColor = None
PP_ALIGN = MSO_ANCHOR = MSO_AUTO_SIZE = MSO_SHAPE = None
CategoryChartData = XL_CHART_TYPE = XL_LEGEND_POSITION = None

# 16:9 canvas, in inches.
SLIDE_W, SLIDE_H, MARGIN = 13.333, 7.5, 0.92


def _ensure_pptx() -> bool:
    global Presentation, Inches, Pt, RGBColor
    global PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE, MSO_SHAPE
    if Presentation is not None:
        return True

    def _imp() -> dict:
        from pptx import Presentation as _P
        from pptx.chart.data import CategoryChartData as _CCD
        from pptx.dml.color import RGBColor as _RGB
        from pptx.enum.chart import XL_CHART_TYPE as _CT
        from pptx.enum.chart import XL_LEGEND_POSITION as _LP
        from pptx.enum.shapes import MSO_SHAPE as _SH
        from pptx.enum.text import MSO_ANCHOR as _AN
        from pptx.enum.text import MSO_AUTO_SIZE as _AS
        from pptx.enum.text import PP_ALIGN as _AL
        from pptx.util import Inches as _I
        from pptx.util import Pt as _Pt

        return {
            "Presentation": _P, "Inches": _I, "Pt": _Pt, "RGBColor": _RGB,
            "PP_ALIGN": _AL, "MSO_ANCHOR": _AN, "MSO_AUTO_SIZE": _AS, "MSO_SHAPE": _SH,
            "CategoryChartData": _CCD, "XL_CHART_TYPE": _CT, "XL_LEGEND_POSITION": _LP,
        }

    from tools.lazy_deps import ensure_and_bind

    return ensure_and_bind("document.pptx", _imp, globals(), prompt=False)


@dataclass(frozen=True)
class Theme:
    bg: str
    ink: str  # titles
    body: str  # body text
    accent: str
    muted: str
    font: str = "Arial"  # widely available; LibreOffice substitutes a metric-compatible face


THEMES: dict[str, Theme] = {
    "light": Theme(bg="FFFFFF", ink="1A1D21", body="3A3F45", accent="2F8F6B", muted="8A9099"),
    "dark": Theme(bg="14171C", ink="F4F5F7", body="C7CCD3", accent="49B488", muted="7B828B"),
}

_ALIGN = {}  # filled after import in _A()
_ANCHOR = {}


def _A():
    """Map align/anchor names to enums (after lazy import)."""
    global _ALIGN, _ANCHOR
    if not _ALIGN:
        _ALIGN = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT}
        _ANCHOR = {"top": MSO_ANCHOR.TOP, "middle": MSO_ANCHOR.MIDDLE, "bottom": MSO_ANCHOR.BOTTOM}


def _rgb(hexstr: str):
    return RGBColor.from_string(hexstr)


def _set_bg(slide, theme: Theme) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = _rgb(theme.bg)


def _accent_rule(slide, theme: Theme, x: float, y: float, w: float = 0.9, h: float = 0.07) -> None:
    rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    rect.fill.solid()
    rect.fill.fore_color.rgb = _rgb(theme.accent)
    rect.line.fill.background()
    rect.shadow.inherit = False


def _box(slide, x, y, w, h, anchor="top"):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE  # shrink-to-fit at render time
    _A()
    tf.vertical_anchor = _ANCHOR[anchor]
    return tf


def _add(tf, text, *, size, color, theme, bold=False, italic=False, align="left",
         bullet=False, space_after=8, first=False):
    _A()
    p = tf.paragraphs[0] if (first and not tf.paragraphs[0].runs) else tf.add_paragraph()
    p.alignment = _ALIGN[align]
    p.space_after = Pt(space_after)
    run = p.add_run()
    run.text = ("•  " + text) if bullet else text  # predictable bullets across renderers
    f = run.font
    f.size = Pt(size)
    f.bold = bold
    f.italic = italic
    f.name = theme.font
    f.color.rgb = _rgb(color)
    return p


def _as_list(v: Any) -> list[str]:
    if v is None:
        return []
    if isinstance(v, str):
        return [v]
    return [str(x) for x in v]


# ── layouts ───────────────────────────────────────────────────────────────────
def _blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])  # truly blank


def lay_title(prs, theme, s):
    sl = _blank(prs)
    _set_bg(sl, theme)
    _accent_rule(sl, theme, MARGIN, 2.55, 1.1, 0.09)
    tf = _box(sl, MARGIN, 2.8, SLIDE_W - 2 * MARGIN, 2.2, anchor="top")
    _add(tf, s.get("title", ""), size=46, color=theme.ink, theme=theme, bold=True, first=True, space_after=10)
    if s.get("subtitle"):
        _add(tf, s["subtitle"], size=22, color=theme.muted, theme=theme)


def lay_section(prs, theme, s):
    sl = _blank(prs)
    _set_bg(sl, theme)
    _accent_rule(sl, theme, MARGIN, 3.35, 1.1, 0.09)
    tf = _box(sl, MARGIN, 3.6, SLIDE_W - 2 * MARGIN, 1.8, anchor="top")
    _add(tf, s.get("title", ""), size=38, color=theme.ink, theme=theme, bold=True, first=True)


def _title_band(sl, theme, title):
    tf = _box(sl, MARGIN, MARGIN, SLIDE_W - 2 * MARGIN, 1.0, anchor="top")
    _add(tf, title or "", size=28, color=theme.ink, theme=theme, bold=True, first=True)
    _accent_rule(sl, theme, MARGIN, MARGIN + 0.92, 0.8, 0.06)


def lay_bullets(prs, theme, s):
    sl = _blank(prs)
    _set_bg(sl, theme)
    _title_band(sl, theme, s.get("title"))
    tf = _box(sl, MARGIN, 2.15, SLIDE_W - 2 * MARGIN, SLIDE_H - 2.15 - MARGIN, anchor="top")
    items = _as_list(s.get("bullets"))
    for i, b in enumerate(items):
        _add(tf, b, size=18, color=theme.body, theme=theme, bullet=True, space_after=12, first=(i == 0))


def _column(sl, theme, x, w, col_title, items):
    tf = _box(sl, x, 2.15, w, SLIDE_H - 2.15 - MARGIN, anchor="top")
    if col_title:
        _add(tf, col_title, size=18, color=theme.accent, theme=theme, bold=True, space_after=10, first=True)
    for i, b in enumerate(_as_list(items)):
        _add(tf, b, size=16, color=theme.body, theme=theme, bullet=True, space_after=10,
             first=(i == 0 and not col_title))


def lay_two_column(prs, theme, s):
    sl = _blank(prs)
    _set_bg(sl, theme)
    _title_band(sl, theme, s.get("title"))
    gut = 0.5
    colw = (SLIDE_W - 2 * MARGIN - gut) / 2
    _column(sl, theme, MARGIN, colw, s.get("left_title"), s.get("left"))
    _column(sl, theme, MARGIN + colw + gut, colw, s.get("right_title"), s.get("right"))


def lay_quote(prs, theme, s):
    sl = _blank(prs)
    _set_bg(sl, theme)
    tf = _box(sl, MARGIN + 0.4, 2.2, SLIDE_W - 2 * MARGIN - 0.8, 3.1, anchor="middle")
    _add(tf, "“" + s.get("quote", "") + "”", size=28, color=theme.ink, theme=theme,
         italic=True, first=True, space_after=14)
    if s.get("attribution"):
        _add(tf, "— " + s["attribution"], size=18, color=theme.muted, theme=theme)


def lay_closing(prs, theme, s):
    sl = _blank(prs)
    _set_bg(sl, theme)
    tf = _box(sl, MARGIN, 2.9, SLIDE_W - 2 * MARGIN, 2.0, anchor="middle")
    _add(tf, s.get("title", ""), size=40, color=theme.ink, theme=theme, bold=True, first=True, space_after=10)
    if s.get("subtitle"):
        _add(tf, s["subtitle"], size=20, color=theme.muted, theme=theme)


def lay_table(prs, theme, s):
    sl = _blank(prs)
    _set_bg(sl, theme)
    _title_band(sl, theme, s.get("title"))
    headers = _as_list(s.get("headers"))
    rows = s.get("rows") or []
    ncols = len(headers) or (len(rows[0]) if rows else 1)
    nrows = (1 if headers else 0) + len(rows)
    if nrows == 0:
        return
    top, h = 2.2, min(SLIDE_H - 2.2 - MARGIN, 0.5 * nrows + 0.4)
    gt = sl.shapes.add_table(nrows, ncols, Inches(MARGIN), Inches(top),
                             Inches(SLIDE_W - 2 * MARGIN), Inches(h))
    tbl = gt.table
    r0 = 0
    if headers:
        for c, head in enumerate(headers[:ncols]):
            cell = tbl.cell(0, c)
            cell.text = str(head)
            for para in cell.text_frame.paragraphs:
                for run in para.runs:
                    run.font.bold = True
                    run.font.size = Pt(14)
                    run.font.name = theme.font
                    run.font.color.rgb = _rgb("FFFFFF")
            cell.fill.solid()
            cell.fill.fore_color.rgb = _rgb(theme.accent)
        r0 = 1
    for ri, row in enumerate(rows):
        for c, val in enumerate(_as_list(row)[:ncols]):
            cell = tbl.cell(r0 + ri, c)
            cell.text = str(val)
            for para in cell.text_frame.paragraphs:
                for run in para.runs:
                    run.font.size = Pt(13)
                    run.font.name = theme.font
                    run.font.color.rgb = _rgb(theme.body)


_CHART_TYPES = {}  # filled lazily


def lay_chart(prs, theme, s):
    sl = _blank(prs)
    _set_bg(sl, theme)
    _title_band(sl, theme, s.get("title"))
    global _CHART_TYPES
    if not _CHART_TYPES:
        _CHART_TYPES = {
            "column": XL_CHART_TYPE.COLUMN_CLUSTERED,
            "bar": XL_CHART_TYPE.BAR_CLUSTERED,
            "line": XL_CHART_TYPE.LINE,
            "pie": XL_CHART_TYPE.PIE,
        }
    cats = _as_list(s.get("categories"))
    series = s.get("series") or {}
    if not cats or not series:
        _add(_box(sl, MARGIN, 2.4, SLIDE_W - 2 * MARGIN, 1.0),
             "No chart data provided.", size=16, color=theme.muted, theme=theme, first=True)
        return
    data = CategoryChartData()
    data.categories = cats
    for name, vals in series.items():
        nums = [float(v) for v in _as_list(vals)]
        data.add_series(str(name), nums)
    ctype = _CHART_TYPES.get(str(s.get("chart_type", "column")), _CHART_TYPES["column"])
    gframe = sl.shapes.add_chart(ctype, Inches(MARGIN), Inches(2.2),
                                 Inches(SLIDE_W - 2 * MARGIN), Inches(SLIDE_H - 2.2 - MARGIN), data)
    chart = gframe.chart
    chart.has_legend = len(series) > 1
    if chart.has_legend:
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.include_in_layout = False


def lay_image(prs, theme, s):
    sl = _blank(prs)
    _set_bg(sl, theme)
    _title_band(sl, theme, s.get("title"))
    path = s.get("image_path") or s.get("image")
    if not path:
        return
    import os

    if not os.path.exists(path):
        _add(_box(sl, MARGIN, 2.4, SLIDE_W - 2 * MARGIN, 1.0),
             f"Image not found: {path}", size=14, color=theme.muted, theme=theme, first=True)
        return
    # Fit within the content area, preserving aspect (python-pptx scales from one dim).
    avail_w = SLIDE_W - 2 * MARGIN
    avail_h = SLIDE_H - 2.2 - MARGIN
    pic = sl.shapes.add_picture(path, Inches(MARGIN), Inches(2.2), width=Inches(avail_w))
    if pic.height > Inches(avail_h):
        pic.width = int(pic.width * (Inches(avail_h) / pic.height))
        pic.height = Inches(avail_h)
    if s.get("caption"):
        cap = _box(sl, MARGIN, SLIDE_H - MARGIN - 0.4, avail_w, 0.4, anchor="top")
        _add(cap, s["caption"], size=12, color=theme.muted, theme=theme, italic=True, first=True)


LAYOUTS = {
    "title": lay_title,
    "section": lay_section,
    "bullets": lay_bullets,
    "content": lay_bullets,  # alias
    "two_column": lay_two_column,
    "quote": lay_quote,
    "closing": lay_closing,
    "table": lay_table,
    "chart": lay_chart,
    "image": lay_image,
}


def build_deck(spec: dict, out_path: str) -> dict:
    """Build a .pptx from a deck spec. Returns {path, slides, warnings}. Never raises on an
    unknown layout — it falls back to a bullets slide and records a warning."""
    if not _ensure_pptx():
        raise RuntimeError("python-pptx unavailable — cannot build the deck")
    if not isinstance(spec, dict) or not isinstance(spec.get("slides"), list):
        raise ValueError("deck spec must be an object with a 'slides' list")
    theme = THEMES.get(str(spec.get("theme", "light")), THEMES["light"])
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)
    warnings: list[str] = []
    for i, s in enumerate(spec["slides"]):
        if not isinstance(s, dict):
            warnings.append(f"slide {i}: not an object, skipped")
            continue
        layout = str(s.get("layout", "bullets"))
        fn = LAYOUTS.get(layout)
        if fn is None:
            warnings.append(f"slide {i}: unknown layout {layout!r} → bullets")
            fn = lay_bullets
        fn(prs, theme, s)
    prs.save(out_path)
    return {"path": out_path, "slides": len(prs.slides._sldIdLst), "warnings": warnings}


# ── markdown/outline → deck spec (the fast path for "make a deck from this") ──
def markdown_to_deck(md: str, *, theme: str = "light", max_bullets: int = 6) -> dict:
    """Turn a markdown outline into a deck spec (flat slide dicts the layouts read).

    First H1 → a title slide; each subsequent #/## heading starts a new bullets
    slide whose bullets are the list items / short lines under it; long sections
    split across slides at `max_bullets`. Deterministic — no model reasoning.
    """
    import re as _re

    lines = md.replace("\r\n", "\n").split("\n")
    if lines and lines[0].strip() == "---":  # strip frontmatter
        for i in range(1, len(lines)):
            if lines[i].strip() == "---":
                lines = lines[i + 1 :]
                break

    def _clean(s: str) -> str:
        s = _re.sub(r"!\[[^\]]*\]\([^)]*\)", "", s)
        s = _re.sub(r"\[([^\]]+)\]\([^)]*\)", r"\1", s)
        s = _re.sub(r"[*_`~]+", "", s)
        return s.strip()

    deck_title = None
    slides: list[dict] = []
    cur_title: str | None = None
    cur_bullets: list[str] = []

    def flush():
        nonlocal cur_bullets
        if cur_title is None and not cur_bullets:
            return
        bl = [b for b in cur_bullets if b]
        if not bl:
            slides.append({"layout": "section", "title": cur_title or ""})
        else:
            for k in range(0, len(bl), max_bullets):
                slides.append({"layout": "bullets", "title": cur_title or "", "bullets": bl[k : k + max_bullets]})
        cur_bullets = []

    for raw in lines:
        line = raw.strip()
        hm = _re.match(r"^(#{1,6})\s+(.*)$", line)
        if hm:
            text = _clean(hm.group(2))
            if len(hm.group(1)) == 1 and deck_title is None and not slides:
                deck_title = text
                slides.append({"layout": "title", "title": text})
                cur_title = None
                continue
            flush()
            cur_title = text
            continue
        lm = _re.match(r"^\s*(?:[-*+]|\d+[.)])\s+(.*)$", raw)
        if lm:
            cur_bullets.append(_clean(lm.group(1)))
            continue
        if line and not line.startswith("|") and not line.startswith("!["):
            # short prose line becomes a bullet (keeps slides readable)
            c = _clean(line)
            if c:
                cur_bullets.append(c)
    flush()
    if not slides:
        slides = [{"layout": "title", "title": deck_title or "Untitled"}]
    spec = {"theme": theme, "slides": slides}
    if deck_title:
        spec["title"] = deck_title
    return spec


# ── in-process tool: build_presentation ────────────────────────────────────────
# A TOOL (not a terminal script): runs in the Robin backend where python-pptx
# auto-installs via lazy_deps. See office_docx.build_document for why.
BUILD_PRESENTATION_SCHEMA = {
    "name": "build_presentation",
    "description": (
        "Build a polished .pptx using designed 16:9 templates that own all layout — never "
        "hand-write python-pptx. FAST PATH: pass `markdown_path` (or `markdown`) to turn an "
        "outline into a deck in one call. Or pass a `spec`: {title, theme:'light'|'dark', "
        "slides:[{layout, ...flat keys}]}; layouts + their keys: title{title,subtitle}, "
        "section{title}, bullets{title,bullets[]}, two_column{title,left_title,left[],right_title,"
        "right[]}, quote{quote,attribution}, closing{title,subtitle}, table{title,headers[],rows[][]}, "
        "chart{title,categories[],series{name:[nums]},chart_type}, image{title,image_path,caption}. "
        "Set verify=true to auto-run the render_check gate. Provide one of markdown_path|markdown|spec."
    ),
    "parameters": {
        "type": "object",
        "properties": {
            "markdown_path": {"type": "string", "description": "Path to a .md outline to convert (fast path)."},
            "markdown": {"type": "string", "description": "Markdown outline text to convert (fast path)."},
            "spec": {"type": "object", "description": "A deck spec (see description) for full control."},
            "out_path": {"type": "string", "description": "Absolute output path ending in .pptx"},
            "theme": {"type": "string", "description": "'light' (default) or 'dark'."},
            "verify": {"type": "boolean", "description": "If true, auto-run render_check and include the verdict."},
        },
        "required": ["out_path"],
    },
}


def _deck_spec_from_args(args: dict) -> dict | str:
    import os
    from pathlib import Path

    theme = str(args.get("theme") or "light")
    md_path = str(args.get("markdown_path") or "").strip()
    md_text = args.get("markdown")
    spec = args.get("spec")
    if md_path:
        try:
            md_text = Path(os.path.expanduser(md_path)).read_text(encoding="utf-8")
        except Exception as e:  # noqa: BLE001
            return f"could not read markdown_path: {e}"
    if isinstance(md_text, str) and md_text.strip():
        return markdown_to_deck(md_text, theme=theme)
    if isinstance(spec, str):
        try:
            spec = json.loads(spec)
        except Exception:
            return "spec must be a JSON object"
    if isinstance(spec, dict):
        spec.setdefault("theme", theme)
        return spec
    return "provide one of markdown_path, markdown, or spec"


async def _build_presentation_handler(args: dict, **_kw) -> str:
    out_path = str(args.get("out_path") or "").strip()
    if not out_path.lower().endswith(".pptx"):
        return json.dumps({"error": "out_path must end in .pptx"})
    spec = _deck_spec_from_args(args)
    if isinstance(spec, str):
        return json.dumps({"error": spec})
    try:
        res = build_deck(spec, out_path)
    except Exception as e:  # noqa: BLE001
        return json.dumps({"error": f"build_presentation failed: {e}"})
    if args.get("verify"):
        try:
            from tools.document_qa import check_document
            res["render_check"] = await check_document(out_path)
        except Exception as e:  # noqa: BLE001
            res["render_check"] = {"error": f"verify failed: {e}"}
    return json.dumps(res)


# Top-level registration so discover_builtin_tools() picks this up (see office_docx).
from tools.registry import registry  # noqa: E402

registry.register(
    name="build_presentation",
    toolset="office",
    schema=BUILD_PRESENTATION_SCHEMA,
    handler=_build_presentation_handler,
    check_fn=lambda: True,  # python-pptx lazy-installs on first use
    emoji="📊",
)
